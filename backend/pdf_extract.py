"""
PDF / PPTX Slide Extractor
Extracts per-slide metadata from a PDF file:
  - page dimensions
  - all text content
  - image count & sizes
  - link/annotation count
  - font info
No AI involved — pure structural extraction via PyMuPDF.
"""

import fitz  # PyMuPDF
import os
import json
import sys
import base64
from io import BytesIO


# ── PPTX extraction (python-pptx) ────────────────────────────────────────────

def _extract_pptx(pptx_bytes: bytes, filename: str, include_thumbnails: bool = True) -> dict:
    """
    Extract per-slide info from a PPTX file using python-pptx.
    Renders thumbnails by converting each slide to a PDF page via fitz.
    """
    from pptx import Presentation
    from pptx.util import Pt
    import tempfile, subprocess, shutil

    prs = Presentation(BytesIO(pptx_bytes))
    slides_out = []

    for i, slide in enumerate(prs.slides, start=1):
        # ── Text ──────────────────────────────────────────────────
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        text_content = "\n".join(texts)

        # ── Images ───────────────────────────────────────────────
        image_count = sum(1 for shape in slide.shapes if shape.shape_type == 13)  # MSO_SHAPE_TYPE.PICTURE

        # ── Links ────────────────────────────────────────────────
        links = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.hyperlink and run.hyperlink.address:
                            links.append({"uri": run.hyperlink.address})

        # ── Thumbnail via fitz (render slide as image) ────────────
        thumbnail_b64 = None
        if include_thumbnails:
            try:
                # Save just this slide's content into a temp PDF via fitz
                # Best effort: use slide dimensions to create a blank page with text
                slide_width = prs.slide_width.pt
                slide_height = prs.slide_height.pt
                doc = fitz.open()
                page = doc.new_page(width=slide_width, height=slide_height)
                # Write text onto the page for a basic visual representation
                if text_content:
                    page.insert_textbox(
                        fitz.Rect(20, 20, slide_width - 20, slide_height - 20),
                        text_content[:500],
                        fontsize=11,
                        color=(0, 0, 0),
                    )
                pix = page.get_pixmap(dpi=72)
                thumbnail_b64 = base64.b64encode(pix.tobytes("png")).decode("ascii")
                doc.close()
            except Exception:
                pass

        slides_out.append({
            "slide": i,
            "source_file": filename,
            "dimensions": {},
            "text_content": text_content,
            "text_block_count": len(texts),
            "images": [],
            "image_count": image_count,
            "links": links,
            "link_count": len(links),
            "fonts": [],
            "thumbnail_b64": thumbnail_b64,
        })

    return {
        "filename": filename,
        "file_size": len(pptx_bytes),
        "total_slides": len(slides_out),
        "metadata": {
            "title": prs.core_properties.title or "",
            "author": prs.core_properties.author or "",
            "subject": prs.core_properties.subject or "",
            "creator": "",
            "producer": "python-pptx",
            "creation_date": str(prs.core_properties.created or ""),
            "modification_date": str(prs.core_properties.modified or ""),
        },
        "slides": slides_out,
    }


def extract_slide_info(page, slide_num: int, include_thumbnail: bool = True) -> dict:
    """
    Extract structured info from a single PDF page (slide).
    Returns a dict with text, images, dimensions, fonts, links, etc.
    """
    # ── Dimensions ────────────────────────────────────────────────
    rect = page.rect
    width_pt = round(rect.width, 2)
    height_pt = round(rect.height, 2)
    width_in = round(rect.width / 72, 2)
    height_in = round(rect.height / 72, 2)

    # ── Text extraction ───────────────────────────────────────────
    text_blocks = page.get_text("blocks")  # list of (x0,y0,x1,y1,text,block_no,block_type)
    text_content = page.get_text("text").strip()

    # ── Images on this page ───────────────────────────────────────
    image_list = page.get_images(full=True)
    images_info = []
    for img in image_list:
        xref = img[0]
        try:
            base_image = page.parent.extract_image(xref)
            images_info.append({
                "xref": xref,
                "width": base_image.get("width", 0),
                "height": base_image.get("height", 0),
                "colorspace": base_image.get("colorspace", 0),
                "bpc": base_image.get("bpc", 0),  # bits per component
                "size_bytes": len(base_image.get("image", b"")),
                "ext": base_image.get("ext", "unknown"),
            })
        except Exception:
            images_info.append({"xref": xref, "error": "could not extract"})

    # ── Links / annotations ───────────────────────────────────────
    links = page.get_links()
    links_info = []
    for link in links:
        link_entry = {"kind": link.get("kind", -1)}
        if link.get("uri"):
            link_entry["uri"] = link["uri"]
        if link.get("page") is not None:
            link_entry["target_page"] = link["page"]
        links_info.append(link_entry)

    # ── Fonts used ────────────────────────────────────────────────
    fonts_raw = page.get_fonts()
    fonts = []
    seen = set()
    for f in fonts_raw:
        name = f[3] if len(f) > 3 else "unknown"
        if name not in seen:
            seen.add(name)
            fonts.append(name)

    # ── Thumbnail (small base64 PNG) ──────────────────────────────
    thumbnail_b64 = None
    if include_thumbnail:
        try:
            pix = page.get_pixmap(dpi=72)  # low-res for thumbnails
            thumbnail_b64 = base64.b64encode(pix.tobytes("png")).decode("ascii")
        except Exception:
            pass

    return {
        "slide": slide_num,
        "dimensions": {
            "width_pt": width_pt,
            "height_pt": height_pt,
            "width_in": width_in,
            "height_in": height_in,
        },
        "text_content": text_content,
        "text_block_count": len(text_blocks),
        "images": images_info,
        "image_count": len(images_info),
        "links": links_info,
        "link_count": len(links_info),
        "fonts": fonts,
        "thumbnail_b64": thumbnail_b64,
    }


def extract_pdf(pdf_path: str = None, pdf_bytes: bytes = None,
                filename: str = "unknown.pdf",
                include_thumbnails: bool = True) -> dict:
    """
    Full PDF/PPTX extraction pipeline.
    Accepts either a file path or raw bytes.
    Returns structured JSON with per-slide info.
    """
    if pdf_path:
        filename = os.path.basename(pdf_path)
        file_size = os.path.getsize(pdf_path)
        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()

    if not pdf_bytes:
        raise ValueError("Provide either pdf_path or pdf_bytes")

    ext = os.path.splitext(filename)[1].lower()
    if ext in (".pptx", ".ppt"):
        return _extract_pptx(pdf_bytes, filename, include_thumbnails)

    file_size = len(pdf_bytes)
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")

    metadata = doc.metadata or {}
    total_pages = len(doc)

    slides = []
    for i, page in enumerate(doc, start=1):
        slide_info = extract_slide_info(page, i, include_thumbnail=include_thumbnails)
        slide_info["source_file"] = filename
        slides.append(slide_info)

    doc.close()

    return {
        "filename": filename,
        "file_size": file_size,
        "total_slides": total_pages,
        "metadata": {
            "title": metadata.get("title", ""),
            "author": metadata.get("author", ""),
            "subject": metadata.get("subject", ""),
            "creator": metadata.get("creator", ""),
            "producer": metadata.get("producer", ""),
            "creation_date": metadata.get("creationDate", ""),
            "modification_date": metadata.get("modDate", ""),
        },
        "slides": slides,
    }


# ── CLI entry point ───────────────────────────────────────────────────────────
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python pdf_extract.py <path/to/file.pdf>")
        sys.exit(1)

    path = sys.argv[1]
    if not os.path.exists(path):
        print(f"Error: '{path}' not found.")
        sys.exit(1)

    print(f"\n📄 Extracting: {path}\n")
    result = extract_pdf(pdf_path=path)

    out_path = path.replace(".pdf", "_extracted.json")
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=2)

    print(f"✅ {result['total_slides']} slides extracted → {out_path}")
    for s in result["slides"]:
        text_preview = s["text_content"][:80].replace("\n", " ")
        print(f"   Slide {s['slide']}: {s['image_count']} imgs, "
              f"{s['link_count']} links, {len(s['text_content'])} chars — \"{text_preview}…\"")
